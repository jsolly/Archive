from pptx import Presentation
from pptx.util import Inches
from io import BytesIO
from urllib import request
from other.my_secrets import get_regression_devext_dbqa_gis

REGRESSION_GIS = get_regression_devext_dbqa_gis()


def get_items_from_folder(
    gis, folder, item_types=None
) -> list:  # folder=None returns the root folder
    items = gis.users.me.items(folder=folder)

    if item_types:
        items = [item for item in items if item.type in item_types]
        return items

    return items


def get_folder_names(gis):
    return gis.users.me.folders


if __name__ == "__main__":
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]

    FOLDERS = get_folder_names(REGRESSION_GIS)[:1]
    print(f"I have {len(FOLDERS)} folders to work on")
    DETAILED_ITEMS = []
    for FOLDER in FOLDERS:
        FOLDER_ITEMS = get_items_from_folder(
            REGRESSION_GIS, FOLDER, item_types=["Dashboard"]
        )
        print(f"working on Folder {FOLDER['title']} with {len(FOLDER_ITEMS)} items")
        for INDEX, FOLDER_ITEM in enumerate(FOLDER_ITEMS):
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            title.text = FOLDER_ITEM.title
            subtitle.text = str(FOLDER_ITEM.snippet) + "\n" + FOLDER_ITEM.id

            left = top = Inches(1)
            URL = "photo.jpg"

            with request.urlopen(URL) as img_obj:
                b = BytesIO(img_obj.read())
                slide.shapes.add_picture(b, left, top)

    prs.save("test.pptx")
